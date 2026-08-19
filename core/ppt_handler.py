import io
import os
import glob
import shutil
import tempfile
import subprocess
import logging

import fitz
from pptx import Presentation

from PySide6.QtGui import (
    QPixmap, QImage, QPainter, QColor
)
from PySide6.QtCore import Qt, QRectF, QBuffer, QIODevice, QThread, Signal
from PySide6.QtWidgets import QGraphicsScene

logger = logging.getLogger(__name__)

EMU_PER_INCH = 914400
DEFAULT_DPI = 96
EMU_PER_PT = 12700

_LO_INSTALL_HINT = (
    "LibreOffice is not installed or not found in PATH.\n\n"
    "Please install LibreOffice to enable PPTX rendering:\n"
    "  Windows: https://www.libreoffice.org/download/download/\n"
    "  macOS:   brew install --cask libreoffice\n"
    "  Linux:   sudo apt install libreoffice-impress"
)


def _emu_to_px(emu, dpi=DEFAULT_DPI):
    return int(emu / EMU_PER_INCH * dpi)


def detect_libreoffice():
    """Return the path to the LibreOffice (or soffice) executable, or None."""
    return shutil.which("libreoffice") or shutil.which("soffice")


class PptRenderThread(QThread):
    """Background thread that converts a PPTX to PDF via LibreOffice headless,
    then renders each page to a QImage via PyMuPDF.

    QImage is reentrant and can be constructed / copied inside a QThread.
    The emitted QImage is a deep copy that owns its own buffer, so it is
    safe to use after the thread exits.  The main thread converts it to a
    QPixmap (which must live on the GUI thread).
    """

    slide_rendered = Signal(int, QImage)       # slide_index, image
    progress_changed = Signal(str)             # human-readable status
    rendering_finished = Signal(int)           # count of slides rendered
    rendering_error = Signal(str)              # error message

    def __init__(self, file_path, slide_count, dpi=200):
        super().__init__()
        self.file_path = file_path
        self.slide_count = slide_count
        self.dpi = dpi
        self._cancelled = False
        self._temp_dir = None

    # ------------------------------------------------------------------ #
    # Public
    # ------------------------------------------------------------------ #
    def cancel(self):
        self._cancelled = True

    # ------------------------------------------------------------------ #
    # QThread.run
    # ------------------------------------------------------------------ #
    def run(self):
        try:
            self._do_render()
        except Exception as exc:
            self.rendering_error.emit(
                f"Unexpected error during rendering:\n{exc}"
            )
            logger.exception("PptRenderThread failed")
        finally:
            self._cleanup()

    # ------------------------------------------------------------------ #
    # Implementation
    # ------------------------------------------------------------------ #
    def _do_render(self):
        lo_path = detect_libreoffice()
        if not lo_path:
            self.rendering_error.emit(_LO_INSTALL_HINT)
            return

        # --- Step 1: PPTX → PDF (LibreOffice headless) ---
        self._temp_dir = tempfile.mkdtemp(prefix="uniboard_ppt_")
        self.progress_changed.emit("Converting presentation to PDF…")

        cmd = [
            lo_path,
            "--headless",
            "--norestore",
            "--nolockcheck",
            "--convert-to", "pdf",
            "--outdir", self._temp_dir,
            "--",
            self.file_path,
        ]
        try:
            result = subprocess.run(
                cmd,
                capture_output=True,
                timeout=300,
                text=True,
            )
        except subprocess.TimeoutExpired:
            self.rendering_error.emit(
                "LibreOffice timed out while converting the presentation.\n"
                "Try opening a smaller file or increase the timeout."
            )
            return

        if result.returncode != 0:
            self.rendering_error.emit(
                f"LibreOffice failed to convert the file:\n"
                f"stderr: {result.stderr.strip()}\n"
                f"stdout: {result.stdout.strip()}"
            )
            return

        # Locate the generated PDF
        pdf_files = glob.glob(os.path.join(self._temp_dir, "*.pdf"))
        if not pdf_files:
            self.rendering_error.emit(
                "LibreOffice did not produce a PDF output.\n"
                f"Output directory: {self._temp_dir}"
            )
            return

        pdf_path = pdf_files[0]

        # --- Step 2: PDF pages → QImage (PyMuPDF) ---
        doc = fitz.open(pdf_path)
        page_count = min(len(doc), self.slide_count)
        zoom = self.dpi / 72.0
        mat = fitz.Matrix(zoom, zoom)

        for i in range(page_count):
            if self._cancelled:
                break
            page = doc[i]
            pix = page.get_pixmap(matrix=mat)

            if pix.samples is None:
                self.progress_changed.emit(
                    f"Rendering slide {i + 1}/{page_count}… (skipped)"
                )
                continue

            # QImage wraps pix.samples without copying.  We must deep-copy
            # immediately so the QImage owns its buffer before pix is freed.
            img = QImage(
                pix.samples,
                pix.width,
                pix.height,
                pix.stride,
                QImage.Format_RGB888,
            )
            img = img.copy()  # deep copy — now safe to emit across threads
            self.slide_rendered.emit(i, img)
            self.progress_changed.emit(f"Rendering slide {i + 1}/{page_count}…")

        doc.close()
        self.rendering_finished.emit(
            min(page_count, self.slide_count) if not self._cancelled else 0
        )

    def _cleanup(self):
        if self._temp_dir and os.path.isdir(self._temp_dir):
            shutil.rmtree(self._temp_dir, ignore_errors=True)
            self._temp_dir = None


class PptHandler:
    """Handles PPTX loading, slide rendering (via LibreOffice + PyMuPDF),
    annotation storage, and annotated-PPTX export."""

    def __init__(self):
        self.prs = None
        self.file_path = None
        self.slides = []
        self.slide_width_emu = 0
        self.slide_height_emu = 0
        self.slide_count = 0
        self._slide_renders: dict[int, QPixmap] = {}
        self._annotations = {}
        self._annotation_scenes = {}

    # ------------------------------------------------------------------ #
    # Lifecycle
    # ------------------------------------------------------------------ #
    def load(self, file_path: str):
        """Open a PPTX file and read slide metadata (count, dimensions).

        This does *not* render any slide images — use ``render_all_slides``
        or ``PptRenderThread`` for that.
        """
        self.file_path = file_path
        self.prs = Presentation(file_path)
        self.slides = list(self.prs.slides)
        self.slide_width_emu = self.prs.slide_width
        self.slide_height_emu = self.prs.slide_height
        self.slide_count = len(self.slides)
        self._slide_renders.clear()
        self._annotations.clear()
        self._annotation_scenes.clear()
        return self.slide_count

    def has_render_engine(self) -> bool:
        """True when LibreOffice is available for rendering."""
        return detect_libreoffice() is not None

    # ------------------------------------------------------------------ #
    # Slide retrieval
    # ------------------------------------------------------------------ #
    def get_slide_pixmap(self, index: int) -> QPixmap | None:
        """Return the pre-rendered pixmap for *index*, or None."""
        return self._slide_renders.get(index)

    def get_slide_size(self, dpi=DEFAULT_DPI):
        if self.prs is None:
            return 0, 0
        w = _emu_to_px(self.slide_width_emu, dpi)
        h = _emu_to_px(self.slide_height_emu, dpi)
        return w, h

    def get_slide_thumbnail(self, index: int, max_width: int = 200) -> QPixmap | None:
        pm = self.get_slide_pixmap(index)
        if pm is None:
            return None
        return pm.scaledToWidth(max_width, Qt.SmoothTransformation)

    # ------------------------------------------------------------------ #
    # Annotation scene storage (kept for API compatibility with PptCanvasView)
    # ------------------------------------------------------------------ #
    def get_annotation_scene(self, index: int) -> QGraphicsScene:
        if index not in self._annotation_scenes:
            self._annotation_scenes[index] = QGraphicsScene()
        return self._annotation_scenes[index]

    def save_current_annotations(self, index: int, scene: QGraphicsScene):
        if index not in self._annotation_scenes:
            self._annotation_scenes[index] = scene
        else:
            old_scene = self._annotation_scenes[index]
            for item in old_scene.items():
                old_scene.removeItem(item)
            for item in scene.items():
                old_scene.addItem(item)

    # ------------------------------------------------------------------ #
    # Annotated PPTX export (python-pptx — unchanged, works correctly)
    # ------------------------------------------------------------------ #
    def export_annotated_pptx(self, output_path: str, annotation_map: dict) -> bool:
        """Export the original PPTX with annotation overlays baked in as images."""
        if self.prs is None:
            return False

        prs_copy = Presentation(self.file_path)

        for i, slide in enumerate(prs_copy.slides):
            scene = annotation_map.get(i)
            if scene is None or not scene.items():
                continue

            w_px = _emu_to_px(self.slide_width_emu)
            h_px = _emu_to_px(self.slide_height_emu)

            annotation_pm = QPixmap(w_px, h_px)
            annotation_pm.fill(QColor(0, 0, 0, 0))
            ap = QPainter(annotation_pm)
            ap.setRenderHint(QPainter.Antialiasing)
            scene.render(ap)
            ap.end()

            img_buffer = QBuffer()
            img_buffer.open(QIODevice.ReadWrite)
            annotation_pm.save(img_buffer, "PNG")
            img_bytes = img_buffer.data().data()
            img_buffer.close()

            python_buffer = io.BytesIO(img_bytes)

            slide.shapes.add_picture(
                python_buffer,
                0, 0,
                self.slide_width_emu,
                self.slide_height_emu,
            )

        prs_copy.save(output_path)
        return True

    # ------------------------------------------------------------------ #
    # Cleanup
    # ------------------------------------------------------------------ #
    def close(self):
        self.prs = None
        self._slide_renders.clear()
        self._annotation_scenes.clear()
